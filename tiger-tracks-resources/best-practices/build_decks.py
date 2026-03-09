#!/usr/bin/env python3
"""Build revised Best Practices decks using the Weekly Meeting Template."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy, os

TEMPLATE = '/Users/cliffsimmons/.openclaw/media/inbound/cabaa504-f38e-4686-bbc0-ed80f9266560.pptx'
OUTDIR = '/Users/cliffsimmons/.openclaw/workspace/tiger-tracks-resources/best-practices'

# Colors from template
INK = RGBColor(0x1A, 0x1A, 0x1A)
CHARCOAL = RGBColor(0x2D, 0x2D, 0x2D)
BONE = RGBColor(0xF5, 0xF0, 0xE8)
TEAL = RGBColor(0x22, 0x9F, 0xA1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)

def make_deck(template_path, slides_data, output_path):
    """Create a deck from template using TITLE_AND_BODY layout (index 2) for content slides
    and TITLE layout (index 0) for title slides."""
    prs = Presentation(template_path)
    
    # Remove all existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    
    for sd in slides_data:
        layout_idx = sd.get('layout', 2)  # Default TITLE_AND_BODY
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        
        # Set title
        if 'title' in sd:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 0:
                    ph.text = sd['title']
                    for para in ph.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.name = 'Georgia'
                            run.font.bold = True
                    break
        
        # Set body/subtitle
        if 'body' in sd:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx in (1, 2) and ph.placeholder_format.type in (2, 4):
                    ph.text = ''
                    tf = ph.text_frame
                    tf.word_wrap = True
                    
                    lines = sd['body'].split('\n')
                    for i, line in enumerate(lines):
                        if i == 0:
                            para = tf.paragraphs[0]
                        else:
                            para = tf.add_paragraph()
                        
                        para.text = line
                        para.font.name = 'Calibri'
                        para.font.size = Pt(14)
                        
                        if line.startswith('•') or line.startswith('-'):
                            para.level = 0
                            para.font.size = Pt(13)
                        elif line.startswith('  •') or line.startswith('  -'):
                            para.level = 1
                            para.font.size = Pt(12)
                    break
        
        # Set subtitle for title slides
        if 'subtitle' in sd:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    ph.text = sd['subtitle']
                    for para in ph.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.name = 'Calibri'
                    break
    
    prs.save(output_path)
    print(f'Saved: {output_path} ({len(prs.slides)} slides)')


# ===================== GOOGLE BEST PRACTICES =====================
google_slides = [
    {'layout': 0, 'title': 'Google Ads\nBest Practices', 'subtitle': 'Tiger Tracks // Updated March 2026'},
    {'layout': 1, 'title': 'Search Campaigns'},
    {'layout': 2, 'title': 'Search: Campaign Structure', 'body': '''• Consolidated campaign structure: Brand + NonBrand
• Split NonBrand only for different conversion events, geo targeting, or budget allocation
• Competitor/Conquest within NonBrand unless separate budget
• Ad Groups organized by theme with combined match types (Exact + Broad)
• No Broad Match for brand names (including competitors)
• Avoid keyword-level Final URLs
• 2026 Update: Leverage Google's auto-created assets for RSAs but always review for accuracy
• 2026 Update: Broad Match now benefits from AI-powered intent matching; pair with Smart Bidding for best results'''},
    {'layout': 2, 'title': 'Search: Ad Copy Best Practices', 'body': '''• Dynamic Keyword Insertion: Use sparingly, NEVER with intentional misspellings
• Thematic RSAs: Every Ad Group should have 3 RSAs with different themes/value props
• Ad Customizers for frequent promotions or seasonal copy
  - Schedule uploads to coincide with promo start/stop dates
  - Great for restricted verticals (avoids triggering policy review)
• 2026 Update: Automatically Created Assets (ACA) now available for all search campaigns
  - Review generated headlines/descriptions weekly
  - Opt out of low-performing auto-assets
• 2026 Update: Conversational experience in Google Ads for campaign creation
• 2026 Update: Image extensions now serve more broadly; ensure brand-safe images are uploaded'''},
    {'layout': 2, 'title': 'Search: New Campaign Launch', 'body': '''• Launch with Max Conversions or Max Conversion Value
• After 30 conversions in 2 weeks, implement tCPA or tROAS
• Initial target within 20% of current efficiency, adjust over following 2 weeks
• 2026 Update: Google now recommends 50+ conversions in 30 days for optimal Smart Bidding
• 2026 Update: New "Acquisition" goals can be layered on top of conversion goals
• Budget capping: Search lost due to rank is fine (bids controlling spend, not budget)
  - If client won't increase spend, lower bids gradually until 0% limited by budget
  - Reducing bids can increase volume 4x+ from same spend'''},
    {'layout': 1, 'title': 'Shopping Campaigns'},
    {'layout': 2, 'title': 'Shopping: Structure & Feed', 'body': '''• Separate Ad Groups for products with drastically different margins
• Bid adjustments at the Ad Group level, not Product Group level
• Start with single campaign; break into Brand/NonBrand when volume warrants
• Apply Brand negative KWs to NonBrand campaign
• GMC Feed Optimization:
  - Optimize product titles (brand, color, size, gender, material)
  - Use Feed Rules to automate title enrichment
  - Keep metafields updated and accurate
  - Custom Labels for reporting and segmentation
  - Exclude Gift Cards from all PLAs
• 2026 Update: Google Merchant Center Next is now standard
  - Automatic product insights and competitive pricing data
  - AI-generated product descriptions available
  - Enhanced free listings integration'''},
    {'layout': 1, 'title': 'Performance Max'},
    {'layout': 2, 'title': 'PMax: Structure & Creative', 'body': '''• Split Asset Groups only for different Creative or Audience needs
• Thematic separation if excess creative; avoid unnecessary segmentation
• Ensure Listing Groups are correctly segmented
• PMax targets should always be higher than Search targets
• Use as many assets as possible per Ad Set
• Allow 3-4 weeks before making creative changes
• 2026 Update: Search term reports and negative KW implementation now available in UI
• 2026 Update: Channel-level reporting now available (Search vs Shopping vs Display vs Video)
• 2026 Update: Brand exclusions expanded; apply to prevent brand cannibalization
• 2026 Update: Asset Group-level conversion goals now supported
• 2026 Update: URL exclusions now available to prevent low-quality placements
• Still recommend PMax Insights report for placement-level spend analysis'''},
    {'layout': 1, 'title': 'Demand Gen'},
    {'layout': 2, 'title': 'Demand Gen: Strategy & Structure', 'body': '''• Better choice than PMax for clients with zero Shopping presence
• One campaign unless different conversion events or budget allocation
• Ad Groups segmented by Audience, not Creative
• Prospecting and Retargeting can share one campaign
• eComm clients: segment by Product Group
• Opt in to all Google channels (including Display)
• Use Optimized Targeting
• 2026 Update: Demand Gen now supports Shorts-only placements
• 2026 Update: Lookalike segments expanded with more granular options
• 2026 Update: Product feeds now support dynamic remarketing within Demand Gen
• 2026 Update: New "Engaged view" conversion counting for YouTube'''},
    {'layout': 2, 'title': 'Demand Gen: Creative Specs', 'body': '''Image Ads:
• Group by theme; always fill all text assets
Video Ads:
• Max 3 videos per ad (one per aspect ratio)
• Choose where videos show:
  - Vertical: Shorts only
  - Square: In-stream, In-feed
  - Horizontal: In-stream, In-feed
• Opt out of Video Enhancements
• Always fill all text assets
Carousel Ads:
• Only without a Product Feed
Catalog Ads:
• Only with a Product Feed
• 2026 Update: AI-powered creative tools now generate video assets from images
• 2026 Update: New interactive ad formats for YouTube Shorts'''},
    {'layout': 1, 'title': 'Account-Wide Best Practices'},
    {'layout': 2, 'title': 'Naming Conventions & Conversion Setup', 'body': '''Naming Conventions:
• Check with clients before renaming
• Intuitive naming, underscores, no spaces
• Ensure UTMs are updated to match

Conversion Setup:
• Clear, appropriate Conversion Action names
• Correct category and Primary/Secondary classification
• 30/3/1 Attribution Window
• Always set up Enhanced Conversions
• Conversion tracking via GTM instead of GA4 where possible
• Connect CRM/offline conversion data
• 2026 Update: Consent mode v2 now mandatory in many markets
• 2026 Update: Enhanced conversions for leads now supports multiple match keys
• 2026 Update: Google Ads Data Manager simplifies first-party data connections'''},
    {'layout': 2, 'title': 'Keyword & Negative KW Management', 'body': '''Keyword & Search Term Scrubs:
• Complete at least biweekly for every campaign
• Don't remove low-spend KWs even without conversions
• Add search terms with appropriate match type to appropriate Ad Group
• Not always best to add to the AG where it converted

Negative Keyword Lists:
• Brand Safety KW List (apply to all campaigns)
• Brand Term List (apply to NonBrand campaigns)
• Regular negative KW mining from search term reports

Assets & Extensions:
• Utilize ALL possible extensions
• Prioritize account-level (Business Name, Logo, Promotions)
• Then campaign-level assets
• 2026 Update: Auto-generated sitelinks and callouts; review weekly'''},
    {'layout': 2, 'title': 'Timeline & Client Expectations', 'body': '''Weeks 1-2: Launch & Learning Phase
• Campaign launch and initial data collection
• Smart Bidding in learning period
• Set expectations: performance will fluctuate
• Do NOT make major changes during this period

Weeks 3-8: Optimization Phase
• Begin keyword/search term scrubs
• Adjust bids based on initial data
• A/B test ad copy and creative
• Implement negative keywords
• Refine audience targeting

Weeks 8+: Scale & Refine
• Campaigns should be stabilized
• Focus on incremental improvements
• Test new campaign types and features
• Quarterly creative refresh
• Review and reallocate budget based on performance'''},
]

# ===================== META BEST PRACTICES =====================
meta_slides = [
    {'layout': 0, 'title': 'Meta Ads\nBest Practices', 'subtitle': 'Tiger Tracks // Updated March 2026'},
    {'layout': 2, 'title': 'Launch Checklist', 'body': '''Don't launch until these pieces are completed:

1. Tracking: Audit by analytics team; CAPI implementation required
2. Creative: Assets that set campaigns up for success (4:5, 1:1, 9:16)
3. Approval: Written approval on strategy and budget
4. QA: New Account Checklist completed

• 2026 Update: Meta Conversions API (CAPI) is now essential for accurate tracking
• 2026 Update: Verify iOS ATT prompt strategy is in place
• 2026 Update: Ensure Meta Pixel + CAPI are running in parallel (deduplication enabled)'''},
    {'layout': 2, 'title': 'Full-Funnel Approach', 'body': '''Concentrate budget at bottom of funnel (conversion campaigns)
Exclude previous purchasers from all stages except Retention

AWARENESS: Prospecting audience
• Reach, LP views, video views objectives
• Broad interest and LAL audiences

PROSPECTING: Warming up brand-aware users
• Advantage+ and 1% LAL Converters
• If niche/small client, use Interest targeting

RETARGETING: Retargeting warm users
• 30-day site visitors (longer for newer/smaller clients)
• Past converters and engagers

RETENTION: Re-engaging past users
• Customer list uploads
• Dynamic product ads

• 2026 Update: Meta now recommends fewer, broader campaigns over hyper-segmented funnels
• 2026 Update: Advantage+ campaigns can now handle full-funnel with audience controls'''},
    {'layout': 2, 'title': 'Naming Conventions', 'body': '''Check with clients before renaming
Intuitive naming, underscores and no spaces
Update UTMs if needed

Campaign: TT_Objective_Funnel-Stage_Conversion-Event_Location_Flighting
Example: TT_ASC_Retargeting_Purchases_Dallas_Evergreen

Ad Set: Audience_Conversion-Event
Example: Broad_Purchases

Ad: CreativeType_CreativeName_Format_Page
Example: Testimonial_Heather_Video_LP

*Flighting & Location are optional'''},
    {'layout': 2, 'title': 'General Best Practices', 'body': '''• Opt in to ALL placements (Advantage+ Placements)
• 4-6 ads per ad set (don't split budget too thin)
• 10+ ads for Advantage+ Shopping campaigns
• 4-6 week creative review and refresh cycle (dependent on budget/frequency)
• Implement CAPI for more accurate measurement
• Don't make calls too early! Ensure enough data supports recommendations

• 2026 Update: Advantage+ campaigns are now Meta's primary recommendation
  - Advantage+ Shopping: Best for eComm with 10+ ad variations
  - Advantage+ App: Standard for app install campaigns
• 2026 Update: Campaign Budget Optimization (CBO) now standard; ad set budgets deprecated for most objectives
• 2026 Update: Engaged-view attribution window now available (video-specific)
• 2026 Update: AI-powered creative tools in Ads Manager for text variations and image generation'''},
    {'layout': 2, 'title': 'Bidding Strategy', 'body': '''Default: Maximize Number of Conversions (Lowest Cost)

Target CPA:
• Use when specific CPA targets are critical
• Limiting but helpful for dramatic seasonality or cultural spikes
• Best with 50+ weekly conversions per ad set

Target ROAS:
• NOT recommended as primary strategy
• Meta conversions are modeled, so data is not 1:1
• Consider only for mature eComm accounts with clean tracking

• 2026 Update: Cost Cap bid strategy now more reliable with larger data pools
• 2026 Update: Minimum ROAS bidding available for catalog campaigns
• 2026 Update: Bid multipliers for Advantage+ Shopping now allow more budget control
• 2026 Update: Performance 5 framework (simplified setup, creative diversification, data integration, quality validation, results analysis)'''},
    {'layout': 2, 'title': 'Audience Targeting', 'body': '''AWARENESS:
• Interest targeting and 5% LAL Converters
• Broad targeting only with significant data and budget (300+ conversions/month)

PROSPECTING:
• Advantage+ and 1% LAL Converters
• Interest targeting if niche/small client
• Exclude audiences to create distinct funnel layers

RETARGETING:
• 30-day site visitors (longer for newer/smaller clients)
• Content engagers and video viewers

RETENTION:
• Past converters and customer lists

• 2026 Update: Advantage+ Audience replaces traditional detailed targeting
  - Meta's AI finds best users; your inputs are "suggestions" not restrictions
• 2026 Update: Advantage Lookalike is now default (auto-expands when beneficial)
• 2026 Update: First-party data (CAPI + customer lists) is the #1 targeting signal'''},
    {'layout': 2, 'title': 'Creative Best Practices', 'body': '''• Matching assets in 4:5, 1:1, and 9:16 ratios for all placements
• Do NOT add faux-CTA buttons to ads
• Variety of ad formats and messaging
• 9:16 placements: add text overlays (copy doesn't pull through)
• Use DPA/Catalog in retargeting
• Clear promotion, offer, or benefit
• 2-4 different creative types per ad set
• Videos under 15 seconds, optimized for sound-off
• Ad copy: not too long, no hashtags or links
• Text under 25% of image area

• 2026 Update: UGC-style and creator content outperforms polished brand creative 2-3x
• 2026 Update: Reels-first creative strategy (vertical video is the primary format)
• 2026 Update: AI creative tools can generate text overlays and crop images automatically
• 2026 Update: Advantage+ Creative now handles dynamic creative optimization; DC deprecated'''},
    {'layout': 2, 'title': 'Ad Specs (March 2026)', 'body': '''Image Ads:
• Recommended: 1,080 x 1,080 (1:1) and 1,080 x 1,920 (9:16)
• Keep top/bottom 15% clear of text in 9:16
• File type: JPG or PNG | Max: 30MB | Min: 600 x 600

Video Ads:
• Recommended: 1,080 x 1,080 (1:1) and 1,080 x 1,920 (9:16)
• File type: MP4, MOV, GIF | Max length: 15s recommended | Max: 4GB

Ad Copy Limits:
• Body text: 125 characters (include on 9:16 overlay)
• Headline: 25 characters
• Link description: 30 characters

• 2026 Update: 9:16 vertical video is now the highest-performing format
• 2026 Update: Carousel ads support up to 20 cards (was 10)
• 2026 Update: Advantage+ Creative automatically generates aspect ratio variations'''},
    {'layout': 2, 'title': 'When to Pause or Refresh Ads', 'body': '''Pause ads when:
• Frequency exceeds 3.0 in prospecting (5.0+ in retargeting)
• CTR drops below 0.5% for 7+ consecutive days
• CPA exceeds target by 50%+ for 2+ weeks with sufficient spend
• Creative fatigue evident (declining engagement + rising CPMs)

Refresh creative when:
• Every 4-6 weeks as standard cadence
• When frequency indicates audience saturation
• Seasonal/promotional changes require new messaging
• Performance plateau despite audience expansion

Do NOT pause when:
• Campaign is still in learning phase (under 50 conversion events)
• Short-term fluctuations during known seasonal periods
• Budget changes made in last 3-5 days (allow stabilization)'''},
    {'layout': 2, 'title': 'Advantage+ Best Practices', 'body': '''Advantage+ Shopping Campaigns (ASC):
• Best for eComm with strong catalog and creative library
• Minimum 10 ad variations; aim for 20+
• Set "existing customer budget cap" (typically 20-30%)
• Use country-level targeting only (let Meta optimize)
• Don't layer additional audience restrictions

Advantage+ Audience:
• Replaces detailed targeting for most objectives
• Your inputs are "suggestions" not hard restrictions
• First-party data (customer lists, site visitors) improves AI signals

Advantage+ Placements:
• Always opt in to all placements
• Meta's AI will allocate budget to best-performing placements
• Manual placement selection only for specific testing needs

Advantage+ Creative:
• Enable for all campaigns
• Handles brightness, contrast, aspect ratio, text enhancements
• Replaces legacy Dynamic Creative Testing (DCT)'''},
    {'layout': 2, 'title': 'Resources', 'body': '''• Meta Ad Specs Reference
• Meta Ads Library (competitor research)
• Learning Period documentation
• Diversified Creative guidance
• Catalog Best Practices
• Tiger Tracks Master Spec Sheet
• Meta Business Help Center: business.facebook.com/help
• Meta Blueprint (training): facebook.com/business/learn'''},
    {'layout': 2, 'title': 'Timeline & Client Expectations', 'body': '''Weeks 1-2: Launch & Learning
• Campaign launch, pixel/CAPI verification
• Learning phase (targeting 50 conversion events)
• Set expectations: costs will be higher, performance unstable
• NO changes during learning period

Weeks 3-8: Optimization
• Creative performance analysis
• Audience refinement
• A/B testing (one variable at a time)
• Frequency monitoring
• First creative refresh

Weeks 8+: Scale & Mature
• Campaigns stabilized
• Expand audiences and test new creative
• Budget scaling (20% increments max)
• Full-funnel analysis and attribution review
• Quarterly strategy reviews'''},
]

# ===================== LINKEDIN BEST PRACTICES =====================
linkedin_slides = [
    {'layout': 0, 'title': 'LinkedIn Ads\nBest Practices', 'subtitle': 'Tiger Tracks // Updated March 2026'},
    {'layout': 2, 'title': 'Why LinkedIn Advertising', 'body': '''• 63 million LinkedIn users are decision-makers
• 80% of users influence business decisions
• 80% of B2B leads via social come from LinkedIn
• 62% of B2B marketers report LinkedIn produces leads
• 20-30% year-over-year ROI uplift for many advertisers
• LinkedIn ads drive ~33% increase in purchase intent
• 13.5% lower CPA with LinkedIn conversion tracking
• 2-5x higher ROAS than other social platforms

• 2026 Update: LinkedIn now reaches 1B+ members globally
• 2026 Update: Revenue Attribution Report links ad interactions to CRM outcomes
• 2026 Update: Predictive Audiences use AI to find high-intent users
• 2026 Update: CTV advertising now available through LinkedIn Campaign Manager'''},
    {'layout': 1, 'title': 'Full-Funnel Strategy'},
    {'layout': 2, 'title': 'Awareness: Raise Brand Visibility', 'body': '''Objectives: Reach, Traffic, Video Views

Audiences:
• Native targeting: job titles, industries, company size, seniority
• Predictive Audiences based on customer lists
• 2026 Update: Company engagement targeting (target employees at companies visiting your site)

Creative Types:
• Static Ads, Carousel, Video
• Ungated Document Ads
• 2026 Update: Click-to-Message ads (direct InMail from ad)

Messaging:
• Introduction to brand and product/service
• Thought leadership positioning

KPIs: Reach, Click Volume, CTR, Video Views, Video Completion Rate'''},
    {'layout': 2, 'title': 'Consideration: Drive Engagement', 'body': '''Objectives: Lead Generation, Engagement

Audiences:
• Retargeting: website visitors, on-platform engagers
• Predictive Audiences
• Target Account Lists (ABM)
• 2026 Update: Engaged Audience segments (interacted with your content)

Creative Types:
• Document Ads (gated and ungated)
• Thought Leadership Ads
• 2026 Update: Newsletter Ads (promote LinkedIn newsletters)
• 2026 Update: Live Event Ads (drive registrations)

Messaging:
• White papers, case studies
• Proven results and testimonials

KPIs: Lead Volume, Conversion Rate, Engagement Rate, Content Downloads'''},
    {'layout': 2, 'title': 'Conversion: Capture High-Intent Leads', 'body': '''Objectives: High-Intent Lead Generation

Audiences:
• Retargeting: content downloaders, high-intent web pages
• Predictive Audiences
• Target Account Lists

Creative Types:
• Message Ads (InMail)
• Lead Form Ads (native forms, no landing page needed)
• Website Conversion Ads
• 2026 Update: Conversions API (CAPI) for server-side tracking

Messaging:
• Drive home benefits and ease of use
• Proven results with statistics
• Clear CTA with urgency

KPIs: Lead Volume, Conversion Rate, Lead Quality, Cost Per Lead, SQL Rate'''},
    {'layout': 2, 'title': 'Reporting & Attribution', 'body': '''LinkedIn Campaign Manager Reporting:
• Revenue Attribution Report links ad interactions to CRM outcomes
• Shows deals won, revenue generated, pipeline created
• Works whether or not leads were captured on LinkedIn

• 2026 Update: Enhanced conversion tracking with LinkedIn Insight Tag + CAPI
• 2026 Update: Multi-touch attribution modeling now available
• 2026 Update: Company Engagement Report shows account-level engagement scores
• 2026 Update: Integration with major CRMs (Salesforce, HubSpot, Dynamics)

Best Practices:
• Always install the LinkedIn Insight Tag
• Set up offline conversion tracking for CRM data
• Use UTM parameters for cross-platform attribution
• Review reports weekly; optimize monthly'''},
    {'layout': 2, 'title': 'Multi-Platform Opportunities', 'body': '''Microsoft Advertising Integration:
• Microsoft Search: Apply LinkedIn audience signals to search campaigns
• Display: Target LinkedIn audiences 1:1 across Microsoft Display Network
• Benefits: Reach B2B audiences while actively searching; bid modifiers for highest-value demos

Connected TV (CTV):
• Target LinkedIn audiences via streaming TV
• 4x more effective than linear TV for B2B targeting
• Partners with Kantar for brand lift studies
• Managed directly in LinkedIn Campaign Manager

• 2026 Update: LinkedIn Audience Network expanded to premium publishers
• 2026 Update: Video ads now serve across LinkedIn + Microsoft ecosystem
• 2026 Update: Programmatic guaranteed deals available for enterprise advertisers'''},
    {'layout': 2, 'title': 'Campaign Structure Best Practices', 'body': '''Naming Conventions:
• Check with clients before renaming
• Format: TT_Platform_Objective_FunnelStage_Audience_Creative
• Example: TT_LI_LeadGen_Consideration_ABM_DocumentAd

Structure:
• Minimum 2 campaigns (Awareness + Conversion)
• Separate campaigns by objective, not audience
• 3-5 ads per campaign minimum
• Budget: Minimum $50/day per campaign for meaningful data

Bidding:
• Start with Maximum Delivery (automated)
• Switch to Manual CPC after 2 weeks of data
• Target CPA bidding available for conversion campaigns
• 2026 Update: Enhanced CPC now default; adjusts bids by up to 45%

Audiences:
• Minimum audience size: 50,000 for Sponsored Content
• Layer 2-3 targeting criteria max (avoid over-narrowing)
• Use Audience Expansion cautiously (monitor quality)
• Exclude current customers from prospecting'''},
    {'layout': 2, 'title': 'Creative Best Practices', 'body': '''Ad Formats & Specs:
• Single Image: 1200 x 627 (1.91:1) or 1080 x 1080 (1:1)
• Video: 16:9 or 1:1, under 30 seconds for awareness, under 90s for consideration
• Carousel: 2-10 cards, 1080 x 1080 each
• Document Ads: PDF up to 100MB, 10+ pages recommended

Copy Guidelines:
• Headline: Under 70 characters (truncates on mobile at 70)
• Introductory text: Under 150 characters (before "see more")
• Speak directly to the audience's role/pain point
• Include data points and specific results
• Strong CTA: Be specific ("Download the Guide" vs "Learn More")

• 2026 Update: Short-form video (under 15s) outperforms on LinkedIn Feed
• 2026 Update: Thought Leadership Ads (boost employee posts) see 2x engagement
• 2026 Update: AI-generated ad copy suggestions now in Campaign Manager
• Test 4+ creative variations per campaign minimum'''},
    {'layout': 2, 'title': 'Timeline & Client Expectations', 'body': '''Weeks 1-2: Launch & Learning
• Campaign setup, Insight Tag verification
• Initial audience and creative testing
• Set expectations: LinkedIn CPCs are higher ($5-12 avg)
• Learning phase: need 15+ conversions per week

Weeks 3-8: Optimization
• Audience performance analysis
• Creative A/B testing
• Lead quality review (SQLs not just MQLs)
• Bid adjustments based on CPA targets
• Demographic reporting review

Weeks 8+: Scale & Refine
• ABM strategy refinement
• Expand to full-funnel (awareness + retargeting)
• Multi-platform integration (Microsoft Ads, CTV)
• Quarterly creative refresh
• Pipeline and revenue attribution review

Key Benchmark Ranges (B2B):
• CPC: $5-12  |  CPL: $50-200  |  CTR: 0.4-0.65%'''},
]

# Build all three decks
make_deck(TEMPLATE, google_slides, os.path.join(OUTDIR, 'TT_Google_Best_Practices_2026.pptx'))
make_deck(TEMPLATE, meta_slides, os.path.join(OUTDIR, 'TT_Meta_Best_Practices_2026.pptx'))
make_deck(TEMPLATE, linkedin_slides, os.path.join(OUTDIR, 'TT_LinkedIn_Best_Practices_2026.pptx'))

print('\nAll 3 decks built successfully!')
